"""
Voice PPT Service
==================
Fetches pre-built voiced PPTs from Supabase Storage,
stamps authenticated MFD's profile onto the narrator strip.

Pre-built PPTs live at: voice-ppt-audio/ppt-templates/{voice}_{lang}.pptx
Per-request: fetch 1 PPT → stamp name/role → return. ~0.5s.
"""

import os
import shutil
import logging
from io import BytesIO
from uuid import UUID
from typing import Optional

import httpx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

from app.config import settings
from app.database import supabase_admin

logger = logging.getLogger(__name__)

# Narrator strip Y position (must match template)
NARRATOR_Y = Inches(4.875)

# Local cache for pre-built PPTs
CACHE_DIR = os.path.join(os.path.dirname(__file__), "..", "..", "ready_ppts")

VALID_VOICES = [
    "sumit", "rahul", "ashutosh", "shubh", "aditya", "sunny",
    "priya", "pooja", "simran", "ishita", "shreya", "kavitha",
]
BUCKET = "voice-ppt-audio"

# Voice metadata
VOICE_META = {
    # Male
    "sumit":    {"gender": "male",   "style": "Professional, balanced warmth"},
    "rahul":    {"gender": "male",   "style": "Composed, builds trust"},
    "ashutosh": {"gender": "male",   "style": "Traditional Hindi narration"},
    "shubh":    {"gender": "male",   "style": "Friendly, versatile default"},
    "aditya":   {"gender": "male",   "style": "Clear, confident delivery"},
    "sunny":    {"gender": "male",   "style": "Energetic, approachable tone"},
    # Female
    "priya":    {"gender": "female", "style": "Upbeat, engaging pitch"},
    "pooja":    {"gender": "female", "style": "Encouraging, guiding"},
    "simran":   {"gender": "female", "style": "Warm conversational"},
    "ishita":   {"gender": "female", "style": "Polished, enterprise-ready"},
    "shreya":   {"gender": "female", "style": "Precise pronunciation, authoritative"},
    "kavitha":  {"gender": "female", "style": "Warm, regional authenticity"},
}


# ──────────────────────────────────────────────
# Template fetching
# ──────────────────────────────────────────────

def _local_path(voice: str, lang: str) -> str:
    return os.path.join(CACHE_DIR, f"{voice}_{lang}.pptx")


def _remote_url(voice: str, lang: str) -> str:
    return (
        f"{settings.SUPABASE_URL}/storage/v1/object/public/"
        f"{BUCKET}/ppt-templates/{voice}_{lang}.pptx"
    )


async def fetch_voiced_template(voice: str, lang: str, output_path: str) -> str:
    """
    Get pre-built voiced PPT. Local cache first → Supabase fallback.
    """
    local = _local_path(voice, lang)

    if os.path.exists(local):
        shutil.copy2(local, output_path)
        logger.info(f"Template {voice}/{lang}: local cache")
        return output_path

    url = _remote_url(voice, lang)
    logger.info(f"Template {voice}/{lang}: fetching from Supabase...")

    async with httpx.AsyncClient(timeout=30) as client:
        resp = await client.get(url)
        if resp.status_code != 200:
            raise Exception(
                f"Pre-built template not found: {voice}_{lang}.pptx "
                f"(HTTP {resp.status_code}). Run pre_build_templates.py first."
            )

        os.makedirs(CACHE_DIR, exist_ok=True)
        with open(local, "wb") as f:
            f.write(resp.content)
        with open(output_path, "wb") as f:
            f.write(resp.content)

        logger.info(f"Template {voice}/{lang}: cached ({len(resp.content)} bytes)")
        return output_path


# ──────────────────────────────────────────────
# MFD Profile
# ──────────────────────────────────────────────

def get_mfd_profile(user_id: UUID) -> dict:
    """
    Fetch MFD profile from mfd_profiles table.
    Uses admin client (service key) since we already authenticated via JWT.
    """
    result = (
        supabase_admin
        .table("mfd_profiles")
        .select("name, phone, google_email")
        .eq("user_id", str(user_id))
        .single()
        .execute()
    )

    if not result.data:
        raise Exception("MFD profile not found. Complete your profile first.")

    return result.data


# ──────────────────────────────────────────────
# Avatar generation
# ──────────────────────────────────────────────

def _create_circular_photo(photo_bytes: bytes, size: int = 200) -> bytes:
    photo = Image.open(BytesIO(photo_bytes)).convert("RGBA")
    photo = photo.resize((size, size), Image.LANCZOS)

    mask = Image.new("L", (size, size), 0)
    ImageDraw.Draw(mask).ellipse([0, 0, size, size], fill=255)

    output = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    output.paste(photo, (0, 0))
    output.putalpha(mask)

    buf = BytesIO()
    output.save(buf, format="PNG")
    return buf.getvalue()


def _create_initials_avatar(name: str, size: int = 200) -> bytes:
    img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)

    draw.ellipse([0, 0, size - 1, size - 1], fill=(13, 148, 136, 255))

    initials = "".join(w[0].upper() for w in name.split()[:2])
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", size // 3)
    except IOError:
        try:
            font = ImageFont.truetype("C:/Windows/Fonts/calibrib.ttf", size // 3)
        except IOError:
            font = ImageFont.load_default()

    bbox = draw.textbbox((0, 0), initials, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.text(((size - tw) // 2, (size - th) // 2 - 4), initials,
              fill=(255, 255, 255, 255), font=font)

    buf = BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


# ──────────────────────────────────────────────
# Stamper — only per-request work
# ──────────────────────────────────────────────

def stamp_distributor(
    ppt_path: str,
    distributor_name: str,
    distributor_role: str = "MF Distributor",
    photo_path: Optional[str] = None,
    output_path: Optional[str] = None,
) -> str:
    """
    Stamp MFD identity onto pre-built voiced PPT.
    Overlays avatar + name + role on narrator strip. ~300ms.
    """
    if output_path is None:
        output_path = ppt_path

    prs = Presentation(ppt_path)

    if photo_path and os.path.exists(photo_path):
        with open(photo_path, "rb") as f:
            avatar_bytes = _create_circular_photo(f.read(), 200)
    else:
        avatar_bytes = _create_initials_avatar(distributor_name, 200)

    for slide in prs.slides:
        # Avatar
        slide.shapes.add_picture(
            BytesIO(avatar_bytes),
            Inches(0.42), NARRATOR_Y + Inches(0.16),
            Inches(0.43), Inches(0.43),
        )

        # Name
        name_box = slide.shapes.add_textbox(
            Inches(1.0), NARRATOR_Y + Inches(0.12),
            Inches(2.5), Inches(0.24),
        )
        tf = name_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = distributor_name
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = "Calibri"

        # Role
        role_box = slide.shapes.add_textbox(
            Inches(1.0), NARRATOR_Y + Inches(0.38),
            Inches(2.5), Inches(0.2),
        )
        tf2 = role_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = distributor_role
        p2.font.size = Pt(7.5)
        p2.font.color.rgb = RGBColor(161, 161, 170)
        p2.font.name = "Calibri"

    prs.save(output_path)
    logger.info(f"Stamped: {distributor_name} → {output_path}")
    return output_path


# ──────────────────────────────────────────────
# Voice options
# ──────────────────────────────────────────────

def get_available_voices() -> dict:
    base_url = f"{settings.SUPABASE_URL}/storage/v1/object/public/{BUCKET}/previews"

    voices = []
    for voice_id in VALID_VOICES:
        meta = VOICE_META.get(voice_id, {"gender": "male", "style": "Natural Indian voice"})
        voices.append({
            "id": voice_id,
            "gender": meta["gender"],
            "style": meta["style"],
            "preview_en": f"{base_url}/{voice_id}_en.mp3",
            "preview_hi": f"{base_url}/{voice_id}_hi.mp3",
        })

    return {
        "voices": voices,
        "languages": {"en": "English", "hi": "Hinglish"},
    }